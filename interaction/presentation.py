from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
import subprocess


# Open an existing presentation
presentation = Presentation("presentation.pptx")

# Get the first slide
slide = presentation.slides[0]

# Add a title to the slide
title = slide.shapes.title
title.text = "My Presentation"

# Add some body text to the slide
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = "This is the body of my presentation."

# Save the presentation to a file
presentation.save("presentation.pptx")




# Open the presentation in PowerPoint
subprocess.run(["powershell", "Start-Process", "-FilePath", "presentation.pptx"])
